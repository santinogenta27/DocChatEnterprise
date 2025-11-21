"""Tool for generating presentations."""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional
from datetime import datetime

from .base_tool import BaseTool, ToolResult

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PresentationTool(BaseTool):
    """Tool for generating PowerPoint presentations from analysis results."""
    
    def get_name(self) -> str:
        return "presentation_generator"
    
    def get_description(self) -> str:
        return "Generate PowerPoint presentations from document analysis results"
    
    def get_keywords(self) -> List[str]:
        return ["presentación", "powerpoint", "ppt", "slides", "diapositivas", "crear presentación"]
    
    def execute(
        self,
        title: str,
        slides: List[Dict[str, str]],
        output_path: Optional[str] = None,
        template: Optional[str] = None,
        **kwargs
    ) -> ToolResult:
        """Generate a presentation."""
        try:
            if output_path is None:
                output_path = self.config.memory_dir / f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Ensure .pptx extension
            if output_path.suffix != '.pptx':
                output_path = output_path.with_suffix('.pptx')
            
            # Generate PowerPoint if available
            if PPTX_AVAILABLE:
                return self._generate_pptx(title, slides, output_path)
            else:
                # Fallback to markdown
                return self._generate_markdown(title, slides, output_path)
        
        except Exception as e:
            return ToolResult(
                success=False,
                data=None,
                message=f"Failed to generate presentation: {str(e)}",
                metadata={"error": str(e)}
            )
    
    def _generate_pptx(self, title: str, slides: List[Dict[str, str]], output_path: Path) -> ToolResult:
        """Generate actual PowerPoint presentation."""
        try:
            # Create presentation
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title_shape.text = title
            subtitle.text = f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            
            # Add content slides
            for slide_data in slides:
                bullet_slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                # Set title
                title_shape = slide.shapes.title
                title_shape.text = slide_data.get('title', 'Slide')
                
                # Set content
                content = slide_data.get('content', '')
                if content:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.text = content
                    
                    # Split into paragraphs if content is long
                    if '\n' in content:
                        lines = content.split('\n')
                        text_frame.text = lines[0]
                        for line in lines[1:]:
                            if line.strip():
                                p = text_frame.add_paragraph()
                                p.text = line.strip()
                                p.level = 0
                
                # Add notes if available
                if slide_data.get('notes'):
                    notes_slide = slide.notes_slide
                    notes_text_frame = notes_slide.notes_text_frame
                    notes_text_frame.text = slide_data['notes']
            
            # Save presentation
            prs.save(str(output_path))
            
            return ToolResult(
                success=True,
                data={"path": str(output_path), "slides": len(slides) + 1},  # +1 for title slide
                message=f"PowerPoint presentation generated: {output_path}",
                metadata={"format": "pptx", "slides": len(slides) + 1}
            )
        
        except Exception as e:
            return ToolResult(
                success=False,
                data=None,
                message=f"Failed to generate PPTX: {str(e)}",
                metadata={"error": str(e)}
            )
    
    def _generate_markdown(self, title: str, slides: List[Dict[str, str]], output_path: Path) -> ToolResult:
        """Generate markdown version as fallback."""
        md_path = output_path.with_suffix('.md')
        content = f"# {title}\n\n"
        content += f"*Generado: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n"
        
        for i, slide in enumerate(slides, 1):
            content += f"## Slide {i}: {slide.get('title', f'Slide {i}')}\n\n"
            content += f"{slide.get('content', '')}\n\n"
            if slide.get('notes'):
                content += f"*Notes: {slide['notes']}*\n\n"
            content += "---\n\n"
        
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return ToolResult(
            success=True,
            data={"path": str(md_path), "slides": len(slides)},
            message=f"Presentation outline generated (Markdown): {md_path}",
            metadata={"format": "markdown", "slides": len(slides)}
        )


